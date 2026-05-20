# ui/chat_view.py
import time
import json
from pathlib import Path
from PySide6.QtWidgets import QWidget, QVBoxLayout, QMessageBox, QFileDialog, QApplication, QListWidgetItem, QAbstractItemView, QSizePolicy
from PySide6.QtCore import Qt, QTimer, QUrl
from PySide6.QtUiTools import QUiLoader
from PySide6.QtGui import QTextCursor

from logic.chat_worker import ChatWorker
from logic.conversation_manager import ConversationManager
from logic.rag_manager import RAGManager
from utils.path_utils import get_resource_path, get_app_settings
from utils.model_config import get_context_limit
from utils.constants import RESPONSE_BUFFER_CHARS

from ui.shared_widgets import MessageData, ChatDisplay

class ChatViewWidget(QWidget):
    def __init__(self, parent_window, llm_client, theme_manager, formatter):
        super().__init__(parent_window)
        self.window = parent_window # Reference to main shell
        self.llm_client = llm_client
        self.theme_manager = theme_manager
        self.formatter = formatter
        self.rag_engine = RAGManager() # Autonomous Local Memory Cache
        
        # Internal State
        self.conversation_manager = ConversationManager()
        self.chat_history = []
        self.current_worker = None
        self.current_response_text = ""
        self.response_start_time = None
        self.is_generating = False
        self.stream_start_position = None
        self.attached_files = []
        self.total_tokens = 0
        self.current_conv_id = None
        self.chat_html_history = [] 
        self._auto_run_sandbox_on_next_response = False

        # Load layout
        loader = QUiLoader()
        ui_file = get_resource_path("ui_designer/chat_mode.ui")
        self.ui = loader.load(str(ui_file), self)
        
        # Set standard layout wrapping the loaded widget
        layout = QVBoxLayout(self)
        layout.setContentsMargins(0, 0, 0, 0)
        layout.addWidget(self.ui)
        
        # 1. Inject custom ChatDisplay just like before
        old_chat = self.ui.chat_display
        self.chat_display = ChatDisplay()
        self.chat_display.link_activated.connect(self.handle_chat_link_action)
        display_layout = self.ui.chat_display_layout
        display_layout.replaceWidget(old_chat, self.chat_display)
        old_chat.setParent(None)
        old_chat.deleteLater()

        # Connect localized widgets
        self.input_field = self.ui.input_field
        self.send_btn = self.ui.send_btn
        self.model_btn = self.ui.model_btn
        self.model_desc_label = self.ui.model_desc_label
        self.auth_btn = self.ui.auth_btn
        self.upload_btn = self.ui.upload_btn
        self.theme_toggle_btn = self.ui.theme_toggle_btn
        self.connection_status_btn = self.ui.connection_status_btn

        # Event filters from parent can be installed or replicated locally
        self.input_field.installEventFilter(self.window) # Let main window handle key hooks

        # Setup Connections
        self.send_btn.clicked.connect(self.handle_send_stop_toggle)
        self.model_btn.clicked.connect(self.window.show_model_popup)
        self.auth_btn.clicked.connect(self.window.handle_auth_button)
        self.upload_btn.clicked.connect(self.handle_upload)
        self.theme_toggle_btn.clicked.connect(self.window.toggle_theme)
        
        # Secure Baseline: Lock inputs down by default on constructor entry
        self.set_chat_enabled(False)
        
        # Sidebar
        self.ui.new_chat_btn.clicked.connect(self.start_new_chat)
        self.ui.delete_all_btn.clicked.connect(self.clear_all_history)
        self.ui.chat_history_list.itemClicked.connect(self.load_selected_history)
        self.ui.chat_history_list.setContextMenuPolicy(Qt.CustomContextMenu)
        self.ui.chat_history_list.customContextMenuRequested.connect(self.show_history_context_menu)
        self.ui.chat_history_list.setSelectionMode(QAbstractItemView.ExtendedSelection)
        
        self.refresh_history_list()
        self.set_send_button_idle()
        self.set_chat_enabled(False)
        
        # Restore Splitter Layout
        self.load_layout_settings()

    def load_layout_settings(self):
        """Restores the visual proportions of the main splitter."""
        settings = get_app_settings()
        state = settings.value("chat_splitter_state")
        if state and hasattr(self.ui, 'main_splitter'):
            self.ui.main_splitter.restoreState(state)

    def save_layout_settings(self):
        """Persists the current visual proportions of the main splitter."""
        if hasattr(self.ui, 'main_splitter'):
            settings = get_app_settings()
            settings.setValue("chat_splitter_state", self.ui.main_splitter.saveState())
            settings.sync()
        
        # --- TOOL INJECTION PHASE ---
        from PySide6.QtWidgets import QCheckBox
        self.web_search_chk = QCheckBox("🌐 Enable Real-Time Web Search")
        # Inject styled grounding toggle directly into visual structure above bottom control bar
        self.web_search_chk.setStyleSheet("""
            QCheckBox {
                color: #0078d4;
                font-weight: bold;
                font-size: 11px;
                margin-left: 15px;
                margin-bottom: 2px;
            }
            QCheckBox::indicator:checked {
                border: 1px solid #0078d4;
                background: #0078d4;
            }
        """)
        # Insert directly above the horizontal bottom toolbar (main_layout node 2)
        self.ui.main_layout.insertWidget(2, self.web_search_chk)

        # Enable dynamic Drag & Drop ingestion matrix
        self.setAcceptDrops(True)
        self.chat_display.setAcceptDrops(False)
        self.input_field.setAcceptDrops(False)

    def shutdown(self):
        """Safely stops background workers to prevent crashes on exit."""
        worker = self.current_worker
        if worker and worker.isRunning():
            worker.terminate()
            worker.wait()
        if hasattr(self, 'vector_sync_thread') and self.vector_sync_thread and self.vector_sync_thread.isRunning():
            # Don't terminate vector thread forcefully to avoid DB corruption
            print("[Shutdown] Waiting for Vector Indexer to finish...")
            self.vector_sync_thread.wait(5000) 

        # Explicitly close the singleton VectorDatabase client connection to avoid interpreter shutdown warnings/race exceptions
        try:
            from logic.vector_db import VectorDatabase
            if VectorDatabase._instance:
                vdb = VectorDatabase.get_instance()
                vdb.close()
        except Exception as e:
            print(f"[Shutdown] Error closing VectorDatabase: {e}") 

    # =========================================================
    # REFACTORED CHAT LOGIC METHODS FROM MAIN_WINDOW
    # =========================================================
    
    def set_chat_enabled(self, enabled):
        self.input_field.setEnabled(enabled)
        self.send_btn.setEnabled(enabled)

    def update_model_ui(self, model_id):
        from logic.model_io import load_all_models
        name = model_id
        desc = ""
        try:
            models = load_all_models()
            for m in models:
                if m['id'] == model_id:
                    name = m.get('name', model_id)
                    desc = m.get('description', '')
                    break
        except Exception as e:
            print(f"Error loading model data: {e}")

        self.model_btn.setText(f"🤖 {name} ▼")
        if desc:
            self.model_desc_label.setText(desc)
            self.model_desc_label.setWordWrap(True)
            self.model_desc_label.setSizePolicy(QSizePolicy.Expanding, QSizePolicy.Minimum)
            self.model_desc_label.setVisible(True)
        else:
            self.model_desc_label.setVisible(False)

    def handle_upload(self):
        """
        ADVANCED FILE LOADER: Expands support to Multiselect, Office Docs, PDFs, 
        and enforces active guard validation against visual triggers.
        """
        if self.is_generating: return
        
        # 1. Ultra-Comprehensive Filter Matrix (Full Office & Data Coverage)
        filter_map = (
            "Total Media Matrix (*.pdf *.docx *.xlsx *.pptx *.odt *.txt *.py *.json *.xml *.ui *.csv *.md *.html *.jpeg *.jpg *.png);"
            "Office Documents (*.pdf *.docx *.xlsx *.pptx *.odt);"
            "Flat Datasheets (*.csv *.txt *.md);"
            "Development & Configs (*.py *.json *.xml *.ui *.html *.yaml *.js);"
            "Graphic Visions (*.jpeg *.jpg *.png)"
        )
        
        # 2. Multi-Select Trigger unlocked
        file_paths, _ = QFileDialog.getOpenFileNames(self, "Attach Assets", "", filter_map)
        
        if not file_paths:
             return
             
        self.process_paths(file_paths)

    def process_paths(self, paths):
        """Direct ingest pipeline supporting both individual files and bulk codebase ingestion."""
        for file_path in paths:
            p = Path(file_path)
            if p.is_dir():
                self.add_system_message(f"📁 **Crawling directory for codebase context:** `{p.name}`...")
                agg_content, count = self.ingest_folder(p)
                if count > 0:
                    # Context Gate: Estimate tokens (char count / 3). 20k tokens ≈ 60,000 chars
                    gate_warning = ""
                    if len(agg_content) > 60000:
                         gate_warning = "\n⚡ *Scale crosses Context Gate (>20k tokens). Local semantic routing ready.*"
                    self.attached_files.append({
                        'name': f"{p.name}/ (Codebase)",
                        'content': agg_content
                    })
                    self.add_system_message(f"✅ **Synthesized {count} source files** from directory: `{p.name}`{gate_warning}")
                else:
                    self.add_system_message(f"⚠️ No supported source code files found in directory: `{p.name}`")
            elif p.is_file():
                self.process_single_file(str(p))

    def process_single_file(self, file_path):
        """Standalone extractor isolating file system IO from interaction handlers."""
        file_ext = Path(file_path).suffix.lower()
        file_name = Path(file_path).name
        
        try:
            # --- TIER 1: VISION GUARD PROTECTION ---
            if file_ext in ['.png', '.jpg', '.jpeg']:
                if not self.llm_client.is_model_vision_capable():
                    QMessageBox.warning(self, "Vision Security Conflict", 
                        f"❌ ATTACHMENT REFUSED: '{file_name}'\n\n"
                        f"The active model lacks a visual recognition stack.\n"
                        "Switch to a valid Vision model to enable visual interaction.")
                    return
                else:
                    # TRUE ACTIVATION: Read binary asset as base64 carrier
                    import base64
                    with open(file_path, 'rb') as img_f:
                        bin_data = base64.b64encode(img_f.read()).decode('utf-8')
                    # Detect refined mime format
                    guessed_mime = f"image/{file_ext.replace('.','').replace('jpg','jpeg')}"
                    self.attached_files.append({
                        'name': file_name, 
                        'content': bin_data, 
                        'type': 'image', 
                        'mime': guessed_mime
                    })
                    # Visual Synergistic Hook: Offer direct-to-sandbox prototype conversion
                    magic_anchor = "<a href='vision_sandbox:run' style='color:#00aaff; text-decoration:none; font-weight:bold;'>🪄 [Compile into Python Prototype]</a>"
                    self.add_system_message(f"🖼️ Loaded Visual: {file_name} &nbsp;&nbsp; {magic_anchor}", allow_html=True)
                    return 

            # --- TIER 2: NATIVE DOCUMENT ENGINE SUITE ---
            elif file_ext == '.pdf':
                import pypdf
                reader = pypdf.PdfReader(file_path)
                content = "\n".join([p.extract_text() for p in reader.pages if p.extract_text()])
                
            elif file_ext == '.docx':
                import docx2txt
                content = docx2txt.process(file_path)

            elif file_ext in ['.xlsx', '.xls']:
                import pandas as pd
                # Fast Excel-to-CSV conversion preserves data topology elegantly for prompt injection
                content = pd.read_excel(file_path).to_csv(index=False)
                
            elif file_ext == '.pptx':
                from pptx import Presentation
                prs = Presentation(file_path)
                slides_txt = []
                for i, slide in enumerate(prs.slides):
                    txt = [f"--- Slide {i+1} ---"]
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            txt.append(shape.text.strip())
                    slides_txt.append("\n".join(txt))
                content = "\n\n".join(slides_txt)
                
            elif file_ext == '.odt':
                from odf import text, teletype
                from odf.opendocument import load
                odtdoc = load(file_path)
                allparas = odtdoc.getElementsByType(text.P)
                content = "\n".join([teletype.extractText(p) for p in allparas])
                
            # --- TIER 3: UNIVERSAL RAW TEXT FALLBACK (JSON, XML, UI, etc.) ---
            else:
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    content = f.read()
                    
            # --- BUNDLE WRITE ---
            if content and content.strip():
                self.attached_files.append({'name': file_name, 'content': content})
                self.add_system_message(f"📎 Synthesized {file_ext.upper()}: {file_name}")
            else:
                self.add_system_message(f"⚠️ Null Content Rejected: {file_name}")
                
        except Exception as e:
            self.add_system_message(f"❌ Parse Crash [{file_name}]: {str(e)}")

    def ingest_folder(self, folder_path):
        """Recursively crawl directory and concatenate supported files in target markdown block formats."""
        folder = Path(folder_path)
        ignored_dirs = {'.git', 'node_modules', '__pycache__', 'venv', '.venv', 'build', 'dist', '.ruby-lsp', '.idea', '.vscode'}
        supported_exts = {'.py', '.js', '.ts', '.html', '.css', '.json', '.xml', '.yaml', '.yml', '.md', '.txt', '.c', '.cpp', '.h', '.hpp', '.java', '.go', '.rs', '.rb', '.php', '.sql', '.sh', '.bat', '.ps1', '.ui'}
        
        agg_content = []
        file_count = 0
        
        def crawl(current_dir, current_depth, max_depth=5):
            nonlocal file_count
            if current_depth > max_depth:
                return
            try:
                for item in current_dir.iterdir():
                    if item.is_dir():
                        if item.name not in ignored_dirs:
                            crawl(item, current_depth + 1, max_depth)
                    elif item.is_file():
                        if item.suffix.lower() in supported_exts:
                            try:
                                with open(item, 'r', encoding='utf-8', errors='ignore') as f:
                                    text = f.read()
                                    if text.strip():
                                        rel_path = item.relative_to(folder)
                                        file_block = (
                                            f"# FILE: {rel_path}\n"
                                            "--- CODE STARTS ---\n"
                                            f"{text}\n"
                                            "--- CODE ENDS ---"
                                        )
                                        agg_content.append(file_block)
                                        file_count += 1
                            except Exception:
                                pass
            except Exception:
                pass
                
        crawl(folder, 1)
        return "\n\n".join(agg_content), file_count

    def dragEnterEvent(self, event):
        """Authorize dynamic drag-and-drop gestures carrying local URIs."""
        if event.mimeData().hasUrls():
            event.acceptProposedAction()

    def dropEvent(self, event):
        """Deconstruct dropped payload and pipe through active ingestion router."""
        paths = [Path(url.toLocalFile()) for url in event.mimeData().urls()]
        valid_paths = [str(p) for p in paths if p.exists()]
        if valid_paths:
             self.process_paths(valid_paths)
             event.acceptProposedAction()

    def clear_attachments(self):
        self.attached_files = []
        self.input_field.setPlaceholderText("")

    def send_message(self, api_response_queue=None, custom_messages=None, custom_temp=None, custom_max_tokens=None, override_prompt=None, api_stream_queue=None):
        # ADAPTIVE PROMPT SOURCE: Favor direct injection logic to avoid UI contamination (Audit ID 039 Fix)
        user_message = override_prompt if override_prompt is not None else self.input_field.toPlainText().strip()
        
        def _api_fail(msg):
            if api_response_queue: api_response_queue.put(f"Error: {msg}")
            if api_stream_queue:
                api_stream_queue.put(f"API_STREAM_ERROR: {msg}")
                api_stream_queue.put(None)

        if not user_message and not self.attached_files:
            _api_fail("Message is empty")
            return

        # Bypass internet check for local providers (Audit ID 037/042 Fix)
        if not self.window.is_connected and not self.llm_client.is_local_provider():
            QMessageBox.warning(self, "No Internet Connection", "Cannot send message.")
            _api_fail("No internet connection")
            return

        if not self.llm_client.current_model:
            self.window.show_model_popup()
            _api_fail("No model selected")
            return

        if not self.llm_client.has_api_key():
            self.window.open_settings()
            _api_fail("Authentication required")
            return

        self.response_start_time = time.perf_counter()  
        
        # Only purge the main textbox if the user manually triggered from GUI
        if override_prompt is None:
             self.input_field.clear()
        
        final_prompt = user_message
        
        # Segregate attachments into functional domains: textual logic vs multimodal vision
        txt_files = [f for f in self.attached_files if f.get('type') != 'image']
        img_files = [f for f in self.attached_files if f.get('type') == 'image']

        if txt_files:
            agg_text = "\n\n".join([f"--- File: {f['name']} ---\n{f['content']}" for f in txt_files])
            
            # DYNAMIC RAG ROUTER: If aggregate text scales past payload safety boundary (15,000 chars),
            # bypass direct prompt dump and redirect stream into fast local vector database automatically.
            if len(agg_text) > 15000:
                self.add_system_message("⚡ **Dataset scales beyond optimal window. Ingesting into local Vector Space...**")
                try:
                    self.rag_engine.ingest_document(agg_text)
                    self.add_system_message("✅ **Local Vector Memory built.** Dynamic semantic retrieval is now active.")
                    final_prompt = user_message if user_message else "Perform deep inference using the semantic vector index built from large local files."
                except Exception as e:
                     self.add_system_message(f"⚠️ **Vector Matrix Collided:** {str(e)}. Falling back to direct context load.")
                     # Emergency Fallback: Just dump it directly despite size, better than losing user data
                     attachment_blocks = [f"--- File: {f['name']} ---\n```\n{f['content']}\n```" for f in txt_files]
                     combined_text = "\n\n".join(attachment_blocks)
                     final_prompt = f"{combined_text}\n\nUser Request: {user_message}" if user_message else f"{combined_text}\n\nPlease review the attached file(s)."
            else:
                # Dataset is small enough for clean, direct context injection
                attachment_blocks = [f"--- File: {f['name']} ---\n```\n{f['content']}\n```" for f in txt_files]
                combined_text = "\n\n".join(attachment_blocks)
                final_prompt = f"{combined_text}\n\nUser Request: {user_message}" if user_message else f"{combined_text}\n\nPlease review the attached file(s)."

        display_msg = user_message if user_message else f"📎 Sent {len(self.attached_files)} asset(s)."
        self.add_user_message(display_msg)

        # Construct Final Payload Carrier (Support Mixed Multimodal Array)
        if img_files:
            # Advanced Carrier: Pack textual directive alongside binary asset frames
            carrier = [{"type": "text", "text": final_prompt}]
            for img in img_files:
                carrier.append({
                    "type": "image",
                    "data": img['content'],
                    "mime": img['mime']
                })
            self.chat_history.append({"role": "user", "content": carrier})
        else:
            # Legacy Carrier: Standard optimized string
            self.chat_history.append({"role": "user", "content": final_prompt})
            
        self.clear_attachments()
        
        # Context Buffer / Compression Check
        TOKEN_BUFFER = 2000
        estimated_new_tokens = len(final_prompt) // 3
        current_total_estimate = self.total_tokens + estimated_new_tokens
        token_limit = get_context_limit(self.llm_client.current_model)
        usage_percent = (current_total_estimate / token_limit) * 100
        
        if (token_limit - current_total_estimate) < TOKEN_BUFFER or usage_percent > 85:
             self.attempt_adaptive_compression(api_response_queue, custom_messages, custom_temp, custom_max_tokens, api_stream_queue)
        else:
             self._continue_send_message(api_response_queue, custom_messages, custom_temp, custom_max_tokens, api_stream_queue)

    def attempt_adaptive_compression(self, api_queue, c_msgs, c_temp, c_max, api_stream_queue=None):
        self.add_system_message("⚡ *Context Optimization Activated...*")
        self.input_field.setEnabled(False)
        self.add_typing_indicator()

        base_history = [m for m in self.chat_history if m.get('role') != 'system']
        if len(base_history) < 3:
             self.remove_typing_indicator()
             self._continue_send_message(api_queue, c_msgs, c_temp, c_max, api_stream_queue)
             return

        pack_count = max(2, int(len(base_history) * 0.60))
        candidates = base_history[:pack_count]
        blueprint = [{"role": "system", "content": "Synthesize the following conversation into ONE highly dense technical paragraph."}]
        for entry in candidates:
             safe_c = self._extract_safe_text(entry.get('content', ''))
             blueprint.append({"role": "user", "content": f"[{entry.get('role', 'user').upper()}]: {safe_c}"})
        blueprint.append({"role": "user", "content": "Execute synthesis."})

        self.compactor_thread = ChatWorker(self.llm_client, blueprint, temperature=0.2, max_tokens=500, parent=self)
        self.compactor_thread.stream = False

        def on_compaction_resolved(text_summary):
             self.remove_typing_indicator()
             pruned = 0
             remaining = []
             for m in self.chat_history:
                  if m.get('role') != 'system' and pruned < pack_count:
                       pruned += 1
                       continue
                  remaining.append(m)
             self.chat_history = remaining
             self.chat_history.insert(0, {"role": "system", "content": f"📊 CONSOLIDATED MEMORY RECAP: {text_summary.strip()}"})
             self.total_tokens = sum(len(m.get('content', '')) // 3 for m in self.chat_history)
             self.add_system_message("✨ *Context Optimized.*")
             self._continue_send_message(api_queue, c_msgs, c_temp, c_max, api_stream_queue)

        def on_compaction_fault(error_string):
             self.remove_typing_indicator()
             self.add_system_message(f"⚠️ *Optimization Bypassed.*")
             self._continue_send_message(api_queue, c_msgs, c_temp, c_max, api_stream_queue)

        self.compactor_thread.response_received.connect(on_compaction_resolved)
        self.compactor_thread.error_occurred.connect(on_compaction_fault)
        self.compactor_thread.start()

    def _continue_send_message(self, api_queue=None, custom_msgs=None, c_temp=None, c_max=None, api_stream_queue=None):
        self.input_field.setEnabled(False)
        self.add_typing_indicator()
        self.current_response_text = "" 
        
        settings = get_app_settings()
        use_defaults = settings.value("gen_use_defaults", "false") == "true"
        
        if use_defaults and c_temp is None and c_max is None:
            a_temp = None
            a_tokens = None
        else:
            a_temp = c_temp if c_temp is not None else float(settings.value("gen_temperature", 0.7))
            a_tokens = c_max if c_max is not None else int(settings.value("gen_max_tokens", 4096))
        
        api_payload = custom_msgs if custom_msgs is not None else self.get_messages_for_api()
        
        # Detect active live search routing
        active_search = None
        if hasattr(self, 'web_search_chk') and self.web_search_chk.isChecked():
             # Extract clean text representation of final payload prompt for searching
             active_search = self._extract_safe_text(self.chat_history[-1].get('content', '')) if self.chat_history else None

        self.current_worker = ChatWorker(
            self.llm_client, 
            api_payload, 
            temperature=a_temp, 
            max_tokens=a_tokens,
            web_search_query=active_search,
            rag_engine=self.rag_engine,
            parent=self
        )    
        self.current_worker.stream_chunk.connect(self.on_stream_chunk)
        self.current_worker.thinking_chunk.connect(self.on_thinking_chunk)
        self.current_worker.response_received.connect(self.on_response_complete)
        self.current_worker.error_occurred.connect(self.on_error)

        if api_queue:
            self.current_worker.response_received.connect(lambda resp: api_queue.put(resp))
            self.current_worker.error_occurred.connect(lambda err: api_queue.put(f"Error: {err}"))
            
        # Hook up continuous streaming bridge back to background Flask callers
        if api_stream_queue:
            self.current_worker.stream_chunk.connect(lambda chunk: api_stream_queue.put(chunk))
            self.current_worker.error_occurred.connect(lambda err: api_stream_queue.put(f"API_STREAM_ERROR:{err}"))
            self.current_worker.finished.connect(lambda: api_stream_queue.put(None)) # Stop sentinel

        self.current_worker.finished.connect(self.on_worker_finished)
        self.current_worker.metrics_received.connect(self.on_metrics_received)
        self.current_worker.start()
        self.set_send_button_generating()

    def add_user_message(self, message: str):
        html = f"<b>You:</b> {self.formatter.escape_html(message)}"
        self.chat_display.append(html)
        self.chat_html_history.append(html)
        self.scroll_to_bottom()

    def add_assistant_message(self, message: str):
        html_content = self.formatter.format_ai_response(message)
        html = f"<b>Assistant:</b><br>{html_content}"
        self.chat_display.append(html)
        self.chat_html_history.append(html)
        self.scroll_to_bottom()

    def add_system_message(self, message: str, allow_html: bool = False):
        if not hasattr(self, 'chat_display'): return
        color = self.theme_manager.get_system_message_color()
        content_text = message if allow_html else self.formatter.escape_html(message)
        html = f"<i style='color: {color};'>ℹ️ {content_text}</i>"
        self.chat_display.append(html)
        self.chat_html_history.append(html)
        self.scroll_to_bottom()

    def add_typing_indicator(self):
        self.chat_display.append("<i>Assistant is typing...</i>")
        self.scroll_to_bottom()

    def remove_typing_indicator(self):
        cursor = self.chat_display.textCursor()
        cursor.movePosition(cursor.MoveOperation.End)
        cursor.movePosition(cursor.MoveOperation.StartOfLine, cursor.MoveMode.KeepAnchor)
        text = cursor.selectedText()
        if "Assistant is typing..." in text:
            cursor.removeSelectedText()
            cursor.deletePreviousChar()

    def on_stream_chunk(self, chunk: str):
        if not self.window.is_connected:
            self.window.is_connected = True
            self.window.update_connection_icon()

        if self.current_response_text == "":
            self.remove_typing_indicator()
            self.chat_display.append("<b>Assistant:</b> ")
            
            # Create a dedicated cursor copy moved to the absolute end to anchor position correctly
            temp_cursor = self.chat_display.textCursor()
            temp_cursor.movePosition(QTextCursor.End)
            self.stream_start_position = temp_cursor.position()

        self.current_response_text += chunk
        cursor = self.chat_display.textCursor()
        cursor.movePosition(cursor.MoveOperation.End)
        self.chat_display.setTextCursor(cursor)
        self.chat_display.insertPlainText(chunk)
        self.scroll_to_bottom()

    def on_thinking_chunk(self, chunk: str):
        text = self.chat_display.toPlainText()
        if "🧠 Thinking..." not in text:
            if self.current_response_text == "":
                self.remove_typing_indicator()
            thinking_color = self.theme_manager.get_thinking_color()
            self.chat_display.append(f"<i style='color: {thinking_color};'>🧠 Thinking...</i>")
            if self.stream_start_position is None:
                # Create a dedicated cursor copy moved to the absolute end to anchor position correctly
                temp_cursor = self.chat_display.textCursor()
                temp_cursor.movePosition(QTextCursor.End)
                self.stream_start_position = temp_cursor.position()

        cursor = self.chat_display.textCursor()
        cursor.movePosition(cursor.MoveOperation.End)
        self.chat_display.setTextCursor(cursor)
        self.chat_display.insertPlainText(chunk)
        self.scroll_to_bottom()

    def on_response_complete(self, full_response: str):
        self.chat_history.append({"role": "assistant", "content": full_response})
        if self.stream_start_position is not None:
            cursor = self.chat_display.textCursor()
            cursor.movePosition(cursor.MoveOperation.End)
            cursor.setPosition(self.stream_start_position, cursor.MoveMode.KeepAnchor)
            cursor.removeSelectedText()
            html_content = self.formatter.format_ai_response(full_response)
            self.chat_display.insertHtml(f"<br>{html_content}<br>")
            
            # Record the assistant block including copy buttons in history log
            copy_buttons = self.theme_manager.get_copy_button_html()
            self.chat_html_history.append(f"<b>Assistant:</b><br>{html_content}<br>{copy_buttons}")
            self.stream_start_position = None
            
            cursor = self.chat_display.textCursor()
            cursor.movePosition(QTextCursor.End)
            cursor.movePosition(QTextCursor.StartOfBlock)
            block = cursor.block()
            block.setUserData(MessageData(full_response))
        else:
            self.chat_display.append("<b>Assistant:</b> ") 
            html = self.formatter.format_ai_response(full_response)
            self.chat_display.insertHtml(html)
            self.chat_html_history.append(f"<b>Assistant:</b><br>{html}<br>{self.theme_manager.get_copy_button_html()}")

        self.auto_save_current_chat()
        self.chat_display.insertHtml(self.theme_manager.get_copy_button_html())

        # 🔄 Vision-to-Sandbox Recursive Execution Loop
        if getattr(self, '_auto_run_sandbox_on_next_response', False):
            self._auto_run_sandbox_on_next_response = False # Reset lock
            try:
                import re, base64
                # Regex extraction targeting Python snippets
                snippet_match = re.search(r"```python\s*\n(.*?)```", full_response, re.DOTALL)
                if not snippet_match:
                    snippet_match = re.search(r"```\s*\n(.*?)```", full_response, re.DOTALL)
                
                if snippet_match:
                    extracted_code = snippet_match.group(1).strip()
                    if extracted_code:
                        self.add_system_message("⚡ **Synergy Triggered: Forwarding generated GUI code directly to Sandbox engine...**")
                        b64_payload = base64.b64encode(extracted_code.encode('utf-8')).decode('utf-8')
                        # Deliver with small visual delay to let typing lock fully release
                        QTimer.singleShot(800, lambda: self.handle_chat_link_action("run_code", b64_payload))
            except Exception as ex:
                 print(f"[Synergy Sandbox Loop] Extraction failed: {ex}")

        self.current_response_text = ""
        if self.response_start_time:
            elapsed = time.perf_counter() - self.response_start_time
            self.chat_display.append(f"<div style='color: #888888; font-size: 12px; text-align: right;'>⏱️ {elapsed:.2f}s</div>")
        self.scroll_to_bottom()

    def on_metrics_received(self, metrics: dict):
        border_color = self.theme_manager.get_metrics_border_color()
        metrics_html = f"<div style='display: flex; gap: 15px; justify-content: flex-end; color: #888888; font-size: 11px; margin-top: 8px; border-top: 1px solid {border_color};'><span>⚡ TTFT: {metrics['ttft']}s</span><span>🚀 Speed: {metrics['tps']} tok/s</span><span>📝 Tokens: {metrics['prompt_tokens']} in / {metrics['completion_tokens']} out</span></div>"
        self.chat_display.append(metrics_html)
        
        # FIX Audit ID 027: Ensure total_tokens represents the ENTIRE conversation history context
        # We sum current prompt context (history + user) and the new AI response.
        self.total_tokens = metrics['prompt_tokens'] + metrics['completion_tokens']
        
        # Safety fallback: If provider returns 0 or single-turn tokens, recalculate from history
        history_estimate = sum(len(m.get('content', '')) // 3 for m in self.chat_history)
        if self.total_tokens < history_estimate:
             self.total_tokens = history_estimate
             
        self.scroll_to_bottom()

    def on_error(self, error_message: str):
        self.remove_typing_indicator()
        err = error_message.lower()
        if "timeout" in err or "connection" in err or "network" in err:
            self.add_system_message("🌐 Network issue. Check connection.")
            self.window.force_disconnected_state()
        else:
            self.add_system_message(f"❌ Error: {error_message}")

        if self.is_generating:
            if self.chat_history and self.chat_history[-1]["role"] == "user": self.chat_history.pop()
            self.is_generating = False
            self.current_worker = None
            self.set_send_button_idle()

    def on_worker_finished(self):
        self.set_send_button_idle()
        self.input_field.setEnabled(True)
        self.input_field.setFocus()
        self.current_worker = None

    def stop_generation(self):
        self.clear_attachments()        
        w = self.current_worker
        self.current_worker = None  
        if w and w.isRunning():
            w.requestInterruption()
            w.quit()
            w.wait(3000)
            if w.isRunning(): w.terminate()
            
        self.remove_typing_indicator()
        self.chat_display.append(f"<br><i style='color: {self.theme_manager.get_terminate_color()};'>⏹️ Terminated</i><br>")
        if self.current_response_text:
            self.chat_history.append({"role": "assistant", "content": self.current_response_text})
        
        self.set_send_button_idle()
        self.input_field.setEnabled(True)
        self.input_field.setFocus()

    def handle_send_stop_toggle(self):
        if self.is_generating: self.stop_generation()
        else: self.send_message()

    def set_send_button_idle(self):
        self.is_generating = False 
        self.send_btn.setText("Send")
        self.send_btn.setStyleSheet("""
            QPushButton { background-color: #0078d4; border-radius: 8px; padding: 12px; color: white; font-weight: bold; }
            QPushButton:hover { background-color: #106ebe; }
            QPushButton:disabled { background-color: #e0e0e0; color: #aaaaaa; border: 1px solid #cccccc; }
        """)

    def set_send_button_generating(self):
        self.is_generating = True
        self.send_btn.setText("Stop")
        self.send_btn.setStyleSheet("""
            QPushButton { background-color: #d32f2f; border-radius: 8px; padding: 12px; color: white; font-weight: bold; }
            QPushButton:hover { background-color: #b71c1c; }
            QPushButton:disabled { background-color: #e0e0e0; color: #aaaaaa; border: 1px solid #cccccc; }
        """)

    def scroll_to_bottom(self):
        scrollbar = self.chat_display.verticalScrollBar()
        scrollbar.setValue(scrollbar.maximum())

    def trigger_vision_to_sandbox_flow(self):
        """Synergistic Contextual Action converting loaded visual assets into auto-executing Python prototypes."""
        visual_assets = [f for f in self.attached_files if f.get('type') == 'image']
        if not visual_assets:
            QMessageBox.information(self, "Asset Missing", "No visual assets currently loaded.\n\nPlease drop or attach an image first!")
            return

        # Set internal recursive loop latch
        self._auto_run_sandbox_on_next_response = True

        directive = (
            "Analyze this visual mock-up/artifact and compile it into a complete, functional, standalone Python desktop application. "
            "Use PySide6 (or Tkinter as fallback) for a modern GUI. Implement active callbacks for mock UI items so they print actions or alter states. "
            "Use HSL tailored curated color palettes and beautiful aesthetics. "
            "Return ONLY a single executable markdown code block containing the complete source code. Zero explanation."
        )
        
        # Dispatch delivery
        self.add_system_message("✨ **Synergistic Protocol Engaged: Compiling Visual into Desktop Prototype...**")
        # Inject directly to stream using sending routing override
        self.send_message(override_prompt=directive)

    def handle_chat_link_action(self, action, payload):
        """Process asynchronous special actions requested from message interactions."""
        import base64
        try:
            decoded = base64.b64decode(payload).decode('utf-8', errors='ignore')
            
            if action == "copy_code":
                QApplication.clipboard().setText(decoded)
                self.add_system_message("📋 **Snippet copied to clipboard.**")
                
            elif action == "vision_sandbox":
                self.trigger_vision_to_sandbox_flow()
                
            elif action == "run_code":
                import sys, tempfile, os
                from PySide6.QtCore import QProcess
                self.add_system_message("⚙️ **Initializing Local Execution Sandbox...**")
                
                # Dump target buffer to temporary runtime artifact
                tmp = tempfile.NamedTemporaryFile(suffix=".py", delete=False)
                tmp.write(decoded.encode('utf-8'))
                tmp.close()
                
                # Instantiate localized computational process bypassing instance variable collisions
                proc = QProcess(self)
                
                def on_sandbox_exit():
                    # Explicitly query targeted process instance local to this invocation's closure
                    out_b = proc.readAllStandardOutput().data()
                    err_b = proc.readAllStandardError().data()
                    stdout = str(out_b, 'utf-8', 'replace').strip()
                    stderr = str(err_b, 'utf-8', 'replace').strip()
                    
                    content = stdout
                    if stderr: content += f"\n\n[RUNTIME ERRORS]:\n{stderr}"
                    if not content: content = "[Execution finished successfully with zero printed output]"
                    
                    self.add_system_message(f"✅ **Execution Output:**\n```bash\n{content}\n```")
                    try: os.unlink(tmp.name) # Clean artifact
                    except: pass
                    
                    # Release object memory instantly on termination back to the event queue
                    proc.deleteLater()
                
                proc.finished.connect(on_sandbox_exit)
                proc.start(sys.executable, [tmp.name])

        except Exception as e:
            self.add_system_message(f"⚠️ Action failed: {str(e)}")

    def clear_chat(self):
        self.chat_display.clear()
        self.chat_history = []
        self.chat_html_history = []
        self.total_tokens = 0
        self.current_response_text = ""
        self.current_conv_id = None
        self.rag_engine.clear() # Purge memory bank for clean session
        self.add_system_message("New conversation started")

    def start_new_chat(self):
        if self.chat_history:
            self.auto_save_current_chat()
        self.clear_chat()
        self.ui.chat_history_list.clearSelection()

    def start_new_chat_without_saving(self):
        self.clear_chat()
        self.ui.chat_history_list.clearSelection()

    def _extract_safe_text(self, content):
        """Isolates pure string text from potentially complex mixed binary lists."""
        if isinstance(content, str): return content
        if isinstance(content, list):
            for item in content:
                if isinstance(item, dict) and item.get('type') == 'text':
                    return item.get('text', '')
        return ""

    def get_messages_for_api(self):
        messages = []
        for msg in self.chat_history:
            if msg.get('role') == 'system':
                if "CONSOLIDATED MEMORY RECAP" in str(msg.get('content', '')):
                    messages.append(msg)
            else: messages.append(msg)
        
        from utils.helpers import get_active_system_instructions
        instructions = get_active_system_instructions()
        if instructions:
            messages.insert(0, {"role": "system", "content": instructions})
        
        def _prep_c(c):
            if isinstance(c, str): return c.strip()
            return c # Forward native complex list payloads untouched
            
        return [{"role": m['role'], "content": _prep_c(m.get('content'))} for m in messages if m.get('content')]

    def auto_save_current_chat(self):
        if not self.chat_history: return
        title = "New Conversation"
        for m in self.chat_history:
            if m['role'] == 'user':
                safe_title = self._extract_safe_text(m.get('content', ''))
                title = safe_title[:30] + "..."
                break
        is_new = (self.current_conv_id is None)
        self.current_conv_id = self.conversation_manager.save_conversation(
            self.chat_history, title=title, conv_id=self.current_conv_id,
            model_id=self.model_btn.text(), messages_html=json.dumps(self.chat_html_history)
        )
        
        # 🧠 Persistent Dense RAG: Queue the latest exchange for background embedding and vector sync
        if len(self.chat_history) >= 2 and self.current_conv_id:
            user_content = None
            assistant_content = None
            # Retrace from reverse to grab latest consecutive user-assistant block
            for i in range(len(self.chat_history) - 1, -1, -1):
                item = self.chat_history[i]
                if item.get('role') == 'assistant' and not assistant_content:
                    assistant_content = self._extract_safe_text(item.get('content', ''))
                elif item.get('role') == 'user' and assistant_content and not user_content:
                    user_content = self._extract_safe_text(item.get('content', ''))
                    break # Block matched

            if user_content and assistant_content:
                try:
                    from workers.vector_indexer_worker import VectorIndexerWorker
                    
                    # 🛡️ INVARIANT 7: Enforce strict single-transaction write serialization for local Qdrant DB
                    if hasattr(self, 'vector_sync_thread') and self.vector_sync_thread and self.vector_sync_thread.isRunning():
                        print("[Dense RAG Ingest] Safeguard triggered: Parallel write suppressed to protect Qdrant lock serialization.")
                    else:
                        # Instantiated thread reference to bypass Python garbage collection
                        self.vector_sync_thread = VectorIndexerWorker(
                            llm_client=self.llm_client,
                            user_text=user_content,
                            assistant_text=assistant_content,
                            conversation_id=self.current_conv_id,
                            model_id=self.model_btn.text(),
                            parent=self
                        )
                        self.vector_sync_thread.start()
                except Exception as vec_e:
                    print(f"[Dense RAG Ingest] Thread setup collision: {vec_e}")

        if is_new: self.refresh_history_list()

    def refresh_history_list(self):
        self.ui.chat_history_list.clear()
        conversations = self.conversation_manager.get_all_conversations()
        for c_id, title, timestamp in conversations:
            item = QListWidgetItem(title)
            item.setData(Qt.ItemDataRole.UserRole, c_id)
            item.setToolTip(f"Saved: {timestamp}")
            self.ui.chat_history_list.addItem(item)

    def load_selected_history(self, item):
        c_id = item.data(Qt.ItemDataRole.UserRole)
        if not c_id: return
        self.current_conv_id = c_id
        try:
            data = self.conversation_manager.load_conversation(c_id)
            if data and "messages" in data:
                self.chat_history = data["messages"]
                self.chat_display.clear()
                cached_html = data.get("messages_html")
                if cached_html:
                    try:
                        self.chat_html_history = json.loads(cached_html)
                        for chunk in self.chat_html_history: self.chat_display.append(chunk)
                    except: pass
                else:
                     for msg in self.chat_history:
                         safe_c = self._extract_safe_text(msg.get('content', ''))
                         if msg['role'] == 'user':
                             self.add_user_message(safe_c)
                         else:
                             self.add_assistant_message(safe_c)
                self.scroll_to_bottom()
        except: pass

    def clear_all_history(self):
        reply = QMessageBox.question(self, "Delete All", "Are you sure you want to delete ALL conversations?", QMessageBox.Yes | QMessageBox.No)
        if reply == QMessageBox.Yes:
            self.conversation_manager.clear_all()
            self.start_new_chat_without_saving()
            self.refresh_history_list()

    def show_history_context_menu(self, pos):
        item = self.ui.chat_history_list.itemAt(pos)
        if not item: return
        from PySide6.QtWidgets import QMenu
        menu = QMenu(self)
        delete_action = menu.addAction("🗑️ Delete")
        if menu.exec(self.ui.chat_history_list.mapToGlobal(pos)) == delete_action:
            self.delete_specific_history(item)

    def delete_specific_history(self, item):
        c_id = item.data(Qt.ItemDataRole.UserRole)
        if not c_id: return
        reply = QMessageBox.question(self, "Delete", f"Delete '{item.text()}'?", QMessageBox.Yes | QMessageBox.No)
        if reply == QMessageBox.Yes:
            self.conversation_manager.delete_conversation(c_id)
            self.refresh_history_list()
            self.start_new_chat_without_saving()

    def save_conversation(self):
        """Export current active memory tree to formatted external JSON"""
        if not self.chat_history:
            QMessageBox.information(self, "Nothing to Save", "No active conversation thread to export.")
            return
        
        file_path, _ = QFileDialog.getSaveFileName(self, "Export Conversation", "", "JSON Files (*.json);;All Files (*)")
        if file_path:
            current_model = self.model_btn.text() if hasattr(self, 'model_btn') else ""
            self.conversation_manager.export_to_json(self.chat_history, file_path, current_model)
            self.add_system_message(f"📂 Active context exported to JSON.")

    def load_conversation(self):
        """Import serialized JSON context tree, hydrating internal history and reconstruction visuals"""
        file_path, _ = QFileDialog.getOpenFileName(self, "Import Conversation", "", "JSON Files (*.json);;All Files (*)")
        if not file_path:
            return
            
        try:
            data = self.conversation_manager.import_from_json(file_path)
            messages = data.get("messages", [])
            
            if not messages:
                QMessageBox.warning(self, "Data Error", "The selected payload does not contain valid context records.")
                return
                
            # 1. Flush active state
            self.clear_chat()
            self.chat_display.clear() # Wiping visual start marker
            
            # 2. Hydrate internal history
            self.chat_history = messages
            
            # 3. Generate Visual Stack reconstructed
            for msg in self.chat_history:
                role = msg.get("role", "user")
                content = msg.get("content", "")
                if not content: continue
                
                # Cleanly isolate text components before pushing to UI widgets
                safe_content = self._extract_safe_text(content)
                if not safe_content: continue
                
                if role == "user":
                    self.add_user_message(safe_content)
                elif role == "assistant":
                    self.add_assistant_message(safe_content)
                elif role == "system":
                    self.add_system_message(safe_content)
                    
            # 4. Revitalize Meta tags
            saved_model = data.get("model_id")
            if saved_model:
                self.llm_client.set_model(saved_model)
                self.update_model_ui(saved_model)
                
            self.add_system_message("✅ Context imported and successfully reconstructed.")
            self.scroll_to_bottom()
            
        except Exception as e:
            QMessageBox.critical(self, "Import Failure", f"Fault encountered rebuilding context payload:\n{str(e)}")
